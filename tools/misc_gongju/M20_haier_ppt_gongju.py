"""
程序简介：处理海尔相关演示文稿生成或整理任务。
主要逻辑：读取所需配置或输入数据，执行本文件负责的处理步骤，并把结果写入本地文件或输出到命令行。
配置说明：涉及企微或飞书凭证时，优先读取 PEIFANG_ENV_PROFILE、WECOM_ENV_PROFILE、FEISHU_ENV_PROFILE 选择公司配置档案；未设置时兼容原来的 .env 变量。
"""

# -*- coding: utf-8 -*-
# v4: 修复导航文字溢出（自动适配），并追加《导航模板页》便于人工复制。
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

INPUT_PPT  = "重庆海尔6237批次杂质问题说明汇报_v4.pptx"
OUTPUT_PPT = "重庆海尔6237批次杂质问题说明汇报_8D合并美化版_v4.pptx"

# —— 颜色 / 字体 —— #
BLUE_MAIN   = RGBColor(0x1F, 0x4E, 0x79)   # 导航描边/未激活文字
BLUE_ACTIVE = RGBColor(0x0F, 0x27, 0x44)   # 导航激活深蓝
GRAY_LIGHT  = RGBColor(0xF2, 0xF4, 0xF7)   # 导航按钮浅灰底
GRAY_LINE   = RGBColor(0xC8, 0xCC, 0xD1)   # 细边
SHADOW_COL  = RGBColor(0xA0, 0xAF, 0xBD)   # 模拟阴影
WHITE       = RGBColor(255, 255, 255)
FONT_CN     = "微软雅黑"
SIZE_NAV    = 14

# —— 导航内容 —— #
NAV_ITEMS = [
    "D1 问题描述",
    "D2 团队组建",
    "D3 初期围堵措施",
    "D4 根本原因分析",
    "D5 纠正措施制定",
    "D6 实施纠正措施",
    "D7 预防再发",
    "D8 效果确认",
]

# 导航栏宽度（并用于右移幅度）
NAV_W = 1.65
# 导航按钮尺寸与间距（略缩 0.02，留出描边/阴影余量，防止视觉“溢出”）
BTN_W, BTN_H = 1.38, 0.53
LEFT_PAD      = 0.18
GAP_Y         = 0.60
# 导航整体起点：标题条高度（确保按钮文字在内容区内）
HEADER_H      = 0.95
TOP_START     = HEADER_H + 0.10
# 投影偏移
SHADOW_DX, SHADOW_DY = 0.04, 0.06

# —— 高亮映射（0-based 页索引）——
ACTIVE_MAP = {
    1: {1},      # 第2页 -> D1
    2: {2, 3},   # 第3页 -> D2 + D3
    3: {4},      # 第4页 -> D4
    4: {5, 6},   # 第5页 -> D5 + D6
    5: {7, 8},   # 第6页 -> D7 + D8
}
# 封面 index=0 不加导航

def add_label_box(slide, x, y, w, h, text, font_color, align=PP_ALIGN.LEFT):
    """带自动缩放适配的文字框"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_CN
    p.font.size = Pt(SIZE_NAV)   # 初始字号
    p.font.color.rgb = font_color
    p.alignment = align
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE   # ✅ 自动缩放以适配按钮
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE           # ✅ 垂直居中
    return tb

def add_nav_button(slide, x, y, label, active=False):
    """绘制一个带投影的圆角按钮；active=True 用深蓝实底白字"""
    # 阴影
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + SHADOW_DX), Inches(y + SHADOW_DY),
        Inches(BTN_W), Inches(BTN_H)
    )
    shadow.fill.solid(); shadow.fill.fore_color.rgb = SHADOW_COL
    shadow.line.color.rgb = SHADOW_COL

    # 主体
    btn = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(BTN_W), Inches(BTN_H)
    )
    btn.line.color.rgb = BLUE_MAIN
    if active:
        btn.fill.solid(); btn.fill.fore_color.rgb = BLUE_ACTIVE
        font_col = WHITE
    else:
        btn.fill.solid(); btn.fill.fore_color.rgb = GRAY_LIGHT
        font_col = BLUE_MAIN

    # 标签（有内边距/自动缩放）
    add_label_box(slide, x + 0.04, y + 0.01, BTN_W - 0.08, BTN_H - 0.02, label, font_col)

def add_left_nav(slide, prs, active_set=None, x0=0.0):
    """添加整栏导航；active_set 是需要高亮的集合（如 {5,6}）"""
    if active_set is None: active_set = set()
    # 导航栏背景：仅覆盖标题条以下区域，避免与深蓝标题重叠
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x0), Inches(HEADER_H), Inches(NAV_W), prs.slide_height - Inches(HEADER_H)
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(245, 247, 250)
    bar.line.color.rgb = GRAY_LINE

    # 8 个按钮
    y = TOP_START
    for i, lab in enumerate(NAV_ITEMS, start=1):
        add_nav_button(slide, x0 + LEFT_PAD, y, lab, active=(i in active_set))
        y += GAP_Y

def process_main():
    prs = Presentation(INPUT_PPT)
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            # 封面不加导航
            continue
        originals = list(slide.shapes)
        active = ACTIVE_MAP.get(idx, set())
        add_left_nav(slide, prs, active_set=active)
        # 右移原内容，避免与导航重叠
        for shp in originals:
            try:
                if shp.width >= Inches(12.6):
                    continue
                shp.left = shp.left + Inches(NAV_W)
            except Exception:
                continue
    return prs

# —— 追加一个《导航模板页》，方便手工复制 —— #
def append_template_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 标题
    title = prs.slides[-1].shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    p.text = "导航模板页（复制整栏到任意页面使用）"
    p.font.name = FONT_CN; p.font.size = Pt(24); p.font.color.rgb = BLUE_ACTIVE

    # 三组导航：无高亮 / 高亮 D5+D6 / 高亮 D7+D8
    add_left_nav(s, prs, active_set=set(),       x0=0.0)
    add_left_nav(s, prs, active_set={5,6},       x0=1.9)
    add_left_nav(s, prs, active_set={7,8},       x0=3.8)

    tip = prs.slides[-1].shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12), Inches(0.6))
    t = tip.text_frame.paragraphs[0]
    t.text = "使用说明：框选整栏（含背景与按钮）→ 右键“组合” → 复制到目标页面 → 如需更换高亮，删除/粘贴另一组即可。"
    t.font.name = FONT_CN; t.font.size = Pt(14); t.font.color.rgb = BLUE_MAIN

def main():
    prs = process_main()
    append_template_slide(prs)
    prs.save(OUTPUT_PPT)
    print("✅ 已生成：", OUTPUT_PPT)

if __name__ == "__main__":
    main()
